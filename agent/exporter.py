# agent/exporter.py
# Type-aware DOCX + PPTX export
# Handles: market | concept | research | mixed

import io
from datetime import date

# ── DOCX ─────────────────────────────────────────────────────────────────────

def generate_docx(query: str, facts: dict, report: str) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc        = Document()
    query_type = facts.get("query_type", "market")

    # ── Styles helpers ────────────────────────────────────────────────────────
    def set_color(run, hex_color):
        r, g, b = int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16)
        run.font.color.rgb = RGBColor(r, g, b)

    def heading(text, level=1, color="1a1a2e"):
        p    = doc.add_heading(text, level=level)
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        for run in p.runs:
            set_color(run, color)
        return p

    def body(text):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(6)
        return p

    def bullet_item(text, color=None):
        p   = doc.add_paragraph(style="List Bullet")
        run = p.add_run(text)
        run.font.size = Pt(11)
        if color:
            set_color(run, color)
        return p

    def section_label(text, color="0066cc"):
        p   = doc.add_paragraph()
        run = p.add_run(text.upper())
        run.bold = True
        run.font.size = Pt(9)
        set_color(run, color)
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after  = Pt(4)
        return p

    def divider():
        doc.add_paragraph("─" * 60)

    # ── Header ────────────────────────────────────────────────────────────────
    title      = doc.add_paragraph()
    title_run  = title.add_run(query)
    title_run.bold      = True
    title_run.font.size = Pt(22)
    set_color(title_run, "0a2463")
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT

    meta      = doc.add_paragraph()
    meta_run  = meta.add_run(f"Generated: {date.today()}  ·  Type: {query_type.upper()}  ·  AI Research Agent")
    meta_run.font.size = Pt(9)
    set_color(meta_run, "888888")

    divider()

    # ── Concept section ───────────────────────────────────────────────────────
    def write_concept(facts):
        if facts.get("definition"):
            section_label("What It Is", "00875a")
            body(facts["definition"])

        if facts.get("simple_analogy"):
            p   = doc.add_paragraph()
            run = p.add_run(f'💡 {facts["simple_analogy"]}')
            run.italic     = True
            run.font.size  = Pt(10)
            set_color(run, "b8860b")

        if facts.get("how_it_works"):
            section_label("How It Works", "0066cc")
            for i, step in enumerate(facts["how_it_works"], 1):
                p   = doc.add_paragraph(style="List Number")
                run = p.add_run(step)
                run.font.size = Pt(11)

        if facts.get("key_components"):
            section_label("Key Components", "5b2d8e")
            for c in facts["key_components"]:
                p    = doc.add_paragraph(style="List Bullet")
                name = p.add_run(f'{c.get("name","")}: ')
                name.bold = True
                name.font.size = Pt(11)
                desc = p.add_run(c.get("description",""))
                desc.font.size = Pt(11)

        if facts.get("use_cases"):
            section_label("Use Cases", "cc6600")
            for u in facts["use_cases"]:
                bullet_item(u, "cc6600")

        if facts.get("common_misconceptions"):
            section_label("Common Misconceptions", "cc0000")
            for m in facts["common_misconceptions"]:
                p   = doc.add_paragraph(style="List Bullet")
                run = p.add_run(f"✗  {m}")
                run.font.size = Pt(11)
                set_color(run, "cc0000")

    # ── Market section ────────────────────────────────────────────────────────
    def write_market(facts):
        if facts.get("companies"):
            section_label("Key Companies", "0066cc")
            for c in facts["companies"]:
                p    = doc.add_paragraph(style="List Bullet")
                name = p.add_run(f'{c.get("name","")}: ')
                name.bold = True
                name.font.size = Pt(11)
                fact = p.add_run(c.get("key_fact", c.get("focus","")))
                fact.font.size = Pt(11)
                ev = c.get("evidence",{})
                if ev.get("source_url"):
                    src = doc.add_paragraph()
                    r   = src.add_run(f'   ↗ {ev.get("source_title","") or ev["source_url"]}')
                    r.font.size = Pt(8)
                    set_color(r, "888888")
                    src.paragraph_format.space_before = Pt(0)

        if facts.get("market_facts"):
            section_label("Market Overview", "006633")
            for f in facts["market_facts"]:
                text = f.get("fact") if isinstance(f, dict) else str(f)
                bullet_item(text or "", "006633")

        if facts.get("challenges"):
            section_label("Challenges", "cc0000")
            for c in facts["challenges"]:
                text = c.get("challenge") if isinstance(c, dict) else str(c)
                bullet_item(text or "", "cc0000")

    # ── Research section ──────────────────────────────────────────────────────
    def write_research(facts):
        if facts.get("key_findings"):
            section_label("Key Findings", "5b2d8e")
            for f in facts["key_findings"]:
                text = f.get("finding") if isinstance(f, dict) else str(f)
                bullet_item(text or "")

        if facts.get("benchmarks"):
            section_label("Benchmarks", "cc6600")
            for b in facts["benchmarks"]:
                p    = doc.add_paragraph(style="List Bullet")
                name = p.add_run(f'{b.get("name","")}: ')
                name.bold = True
                res  = p.add_run(b.get("result",""))
                res.font.size = Pt(11)

        if facts.get("methodologies"):
            section_label("Methodologies", "0066cc")
            for m in facts["methodologies"]:
                bullet_item(m)

        if facts.get("open_problems"):
            section_label("Open Problems", "cc0000")
            for p in facts["open_problems"]:
                bullet_item(p, "cc0000")

        if facts.get("notable_authors"):
            section_label("Notable Authors", "888888")
            body(", ".join(facts["notable_authors"]))

    # ── Write based on type ───────────────────────────────────────────────────
    if query_type == "concept":
        write_concept(facts)
    elif query_type == "market":
        write_market(facts)
    elif query_type == "research":
        write_research(facts)
    elif query_type == "mixed":
        write_concept(facts)
        divider()
        write_market(facts)
    else:
        write_market(facts)

    # ── Full report ───────────────────────────────────────────────────────────
    if report:
        divider()
        section_label("Full Research Report", "0a2463")
        for line in report.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("**") and line.endswith("**"):
                p   = doc.add_paragraph()
                run = p.add_run(line.strip("*"))
                run.bold = True
                run.font.size = Pt(12)
            elif line.startswith("###"):
                heading(line.replace("###","").strip(), level=3)
            else:
                body(line)

    # ── Footer ────────────────────────────────────────────────────────────────
    divider()
    footer      = doc.add_paragraph()
    footer_run  = footer.add_run(f"AI Research Agent  ·  {date.today()}  ·  Powered by OpenAI + Tavily")
    footer_run.font.size = Pt(8)
    set_color(footer_run, "aaaaaa")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def generate_pptx(query: str, facts: dict, report: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    query_type = facts.get("query_type", "market")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Color palette ─────────────────────────────────────────────────────────
    BG      = RGBColor(0x08, 0x0b, 0x0f)
    SURFACE = RGBColor(0x0e, 0x13, 0x18)
    ACCENT  = RGBColor(0x00, 0xe5, 0xff)
    GREEN   = RGBColor(0x00, 0xff, 0x88)
    YELLOW  = RGBColor(0xff, 0xd1, 0x66)
    RED     = RGBColor(0xff, 0x4d, 0x6d)
    PURPLE  = RGBColor(0xc0, 0x84, 0xfc)
    WHITE   = RGBColor(0xff, 0xff, 0xff)
    MUTED   = RGBColor(0x4a, 0x60, 0x70)
    PAPER   = RGBColor(0x81, 0x8c, 0xf8)

    blank_layout = prs.slide_layouts[6]

    def new_slide():
        slide = prs.slides.add_slide(blank_layout)
        bg    = slide.background.fill
        bg.solid(); bg.fore_color.rgb = BG
        return slide

    def add_textbox(slide, text, x, y, w, h,
                    size=18, bold=False, color=WHITE,
                    align=PP_ALIGN.LEFT, wrap=True):
        tb   = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf   = tb.text_frame
        tf.word_wrap = wrap
        p    = tf.paragraphs[0]
        p.alignment = align
        run  = p.add_run()
        run.text      = str(text)
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        return tb

    def add_rect(slide, x, y, w, h, color):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_label(slide, text, x, y, color=ACCENT):
        add_textbox(slide, text.upper(), x, y, 12, 0.3, size=7, bold=True, color=color)
        add_rect(slide, x, y+0.28, 12, 0.02, color)

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = new_slide()
    add_rect(slide, 0, 0, 13.33, 0.06, ACCENT)
    add_rect(slide, 0, 6, 13.33, 0.06, SURFACE)
    add_textbox(slide, "AI RESEARCH AGENT", 0.6, 0.3, 12, 0.4,
                size=8, color=MUTED, bold=True)
    add_textbox(slide, query, 0.6, 0.8, 12, 2.0,
                size=28, bold=True, color=WHITE)

    # Type badge
    type_colors = {"concept":"concept","market":"market","research":"research","mixed":"mixed"}
    badge_colors = {"concept": GREEN, "market": ACCENT, "research": PAPER, "mixed": YELLOW}
    badge_color  = badge_colors.get(query_type, ACCENT)
    add_rect(slide, 0.6, 3.1, 2.5, 0.4, SURFACE)
    add_textbox(slide, f"TYPE: {query_type.upper()}", 0.65, 3.15, 2.4, 0.3,
                size=9, bold=True, color=badge_color)

    # Stats
    stats = []
    if query_type in ("market","mixed"):
        stats.append(f"{len(facts.get('companies',[]))} Companies")
    if query_type in ("concept","mixed"):
        stats.append(f"{len(facts.get('how_it_works',[]))} Steps")
        stats.append(f"{len(facts.get('key_components',[]))} Components")
    if query_type == "research":
        stats.append(f"{len(facts.get('key_findings',[]))} Findings")
        stats.append(f"{len(facts.get('benchmarks',[]))} Benchmarks")

    for i, stat in enumerate(stats[:4]):
        x = 0.6 + i * 3.0
        add_rect(slide, x, 3.7, 2.6, 1.0, SURFACE)
        parts = stat.split(" ", 1)
        add_textbox(slide, parts[0], x+0.15, 3.8, 2.3, 0.5,
                    size=28, bold=True, color=badge_color)
        add_textbox(slide, parts[1] if len(parts)>1 else "", x+0.15, 4.2, 2.3, 0.3,
                    size=9, color=MUTED)

    add_textbox(slide, f"Generated: {date.today()}  ·  Powered by OpenAI + Tavily",
                0.6, 6.8, 12, 0.3, size=8, color=MUTED)

    # ── Concept slides ────────────────────────────────────────────────────────
    def concept_slides():
        # Slide: Definition + How It Works
        if facts.get("definition") or facts.get("how_it_works"):
            s = new_slide()
            add_label(s, "Concept Explanation", 0.5, 0.2, GREEN)

            if facts.get("definition"):
                add_textbox(s, facts["definition"], 0.5, 0.6, 12.3, 1.8,
                            size=13, color=RGBColor(0xc8,0xd8,0xe4))

            if facts.get("simple_analogy"):
                add_rect(s, 0.5, 2.5, 12.3, 0.6, RGBColor(0x1a,0x14,0x00))
                add_textbox(s, f'💡 {facts["simple_analogy"]}', 0.65, 2.55, 12, 0.5,
                            size=10, color=YELLOW)

            steps = facts.get("how_it_works", [])[:5]
            if steps:
                add_label(s, "How It Works", 0.5, 3.3, ACCENT)
                for i, step in enumerate(steps):
                    y = 3.7 + i * 0.55
                    add_rect(s, 0.5, y, 0.35, 0.35, SURFACE)
                    add_textbox(s, str(i+1), 0.55, y+0.02, 0.25, 0.3,
                                size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
                    add_textbox(s, step, 1.0, y, 11.8, 0.45, size=10,
                                color=RGBColor(0xc8,0xd8,0xe4))

        # Slide: Key Components
        components = facts.get("key_components", [])
        if components:
            s = new_slide()
            add_label(s, "Key Components", 0.5, 0.2, PURPLE)
            cols = 3
            for i, comp in enumerate(components[:6]):
                col = i % cols
                row = i // cols
                x   = 0.5 + col * 4.2
                y   = 0.7 + row * 2.8
                add_rect(s, x, y, 3.9, 2.4, SURFACE)
                add_rect(s, x, y, 3.9, 0.06, PURPLE)
                add_textbox(s, comp.get("name",""), x+0.15, y+0.15, 3.6, 0.5,
                            size=13, bold=True, color=WHITE)
                add_textbox(s, comp.get("description",""), x+0.15, y+0.65, 3.6, 1.6,
                            size=9, color=RGBColor(0x4a,0x60,0x70))

        # Slide: Use Cases + Misconceptions
        use_cases = facts.get("use_cases", [])
        misconceptions = facts.get("common_misconceptions", [])
        if use_cases or misconceptions:
            s = new_slide()
            if use_cases:
                add_label(s, "Use Cases", 0.5, 0.2, YELLOW)
                for i, u in enumerate(use_cases[:6]):
                    y = 0.65 + i * 0.6
                    add_rect(s, 0.5, y+0.1, 0.08, 0.25, YELLOW)
                    add_textbox(s, u, 0.75, y, 5.5, 0.55, size=11,
                                color=RGBColor(0xc8,0xd8,0xe4))
            if misconceptions:
                add_label(s, "Common Misconceptions", 7.0, 0.2, RED)
                for i, m in enumerate(misconceptions[:4]):
                    y = 0.65 + i * 0.9
                    add_rect(s, 7.0, y+0.05, 5.8, 0.7, RGBColor(0x1a,0x05,0x08))
                    add_textbox(s, f"✗  {m}", 7.1, y+0.1, 5.6, 0.55, size=10, color=RED)

    # ── Market slides ─────────────────────────────────────────────────────────
    def market_slides():
        companies = facts.get("companies", [])
        if companies:
            per_slide = 4
            for page_start in range(0, min(len(companies), 8), per_slide):
                s     = new_slide()
                chunk = companies[page_start:page_start+per_slide]
                label = "Key Companies" if page_start==0 else f"Companies (cont.)"
                add_label(s, label, 0.5, 0.2)
                cols = 2
                for i, c in enumerate(chunk):
                    col = i % cols
                    row = i // cols
                    x   = 0.5 + col * 6.3
                    y   = 0.7 + row * 3.0
                    add_rect(s, x, y, 6.0, 2.6, SURFACE)
                    add_rect(s, x, y, 6.0, 0.06, ACCENT)
                    add_textbox(s, c.get("name",""), x+0.2, y+0.15, 5.6, 0.5,
                                size=16, bold=True, color=WHITE)
                    add_textbox(s, c.get("focus",""), x+0.2, y+0.65, 5.6, 0.4,
                                size=9, color=MUTED)
                    add_textbox(s, c.get("key_fact",""), x+0.2, y+1.05, 5.6, 1.0,
                                size=10, color=GREEN)
                    ev = c.get("evidence",{})
                    if ev.get("confidence"):
                        conf_colors = {"high":GREEN,"medium":YELLOW,"low":RED}
                        cc = conf_colors.get(ev["confidence"], MUTED)
                        add_textbox(s, ev["confidence"].upper(), x+0.2, y+2.1, 1.0, 0.3,
                                    size=7, bold=True, color=cc)

        market_facts = facts.get("market_facts", [])
        challenges   = facts.get("challenges", [])
        if market_facts or challenges:
            s = new_slide()
            if market_facts:
                add_label(s, "Market Overview", 0.5, 0.2, GREEN)
                for i, f in enumerate(market_facts[:5]):
                    text = f.get("fact") if isinstance(f,dict) else str(f)
                    y    = 0.65 + i * 0.85
                    add_rect(s, 0.5, y+0.15, 0.08, 0.35, GREEN)
                    add_textbox(s, text or "", 0.75, y, 5.5, 0.8,
                                size=11, color=RGBColor(0xc8,0xd8,0xe4))
            if challenges:
                add_label(s, "Challenges", 7.0, 0.2, RED)
                for i, c in enumerate(challenges[:5]):
                    text = c.get("challenge") if isinstance(c,dict) else str(c)
                    y    = 0.65 + i * 0.85
                    add_rect(s, 7.0, y+0.1, 5.8, 0.7, RGBColor(0x1a,0x05,0x08))
                    add_textbox(s, text or "", 7.1, y+0.12, 5.6, 0.6,
                                size=10, color=RED)

    # ── Research slides ───────────────────────────────────────────────────────
    def research_slides():
        findings = facts.get("key_findings", [])
        if findings:
            s = new_slide()
            add_label(s, "Key Findings", 0.5, 0.2, PAPER)
            for i, f in enumerate(findings[:6]):
                text = f.get("finding") if isinstance(f,dict) else str(f)
                y    = 0.65 + i * 0.9
                add_rect(s, 0.5, y+0.15, 0.08, 0.35, PAPER)
                add_textbox(s, text or "", 0.75, y, 11.8, 0.8,
                            size=11, color=RGBColor(0xc8,0xd8,0xe4))

        benchmarks = facts.get("benchmarks", [])
        if benchmarks:
            s = new_slide()
            add_label(s, "Benchmarks", 0.5, 0.2, YELLOW)
            for i, b in enumerate(benchmarks[:6]):
                col = i % 3
                row = i // 3
                x   = 0.5 + col * 4.1
                y   = 0.7 + row * 2.5
                add_rect(s, x, y, 3.8, 2.2, SURFACE)
                add_textbox(s, b.get("name",""), x+0.15, y+0.15, 3.5, 0.5,
                            size=11, bold=True, color=YELLOW)
                add_textbox(s, b.get("result",""), x+0.15, y+0.65, 3.5, 1.3,
                            size=14, bold=True, color=WHITE)

    # ── Conclusion slide ──────────────────────────────────────────────────────
    def conclusion_slide():
        s = new_slide()

        # ── Thank you / closing slide ─────────────────────────────────────────
        # Big centered title
        add_rect(s, 0, 0, 13.33, 7.5, BG)
        add_rect(s, 0, 0, 13.33, 0.06, badge_color)
        add_rect(s, 0, 7.44, 13.33, 0.06, badge_color)

        # Topic title
        add_textbox(s, query, 0.6, 0.5, 12.1, 1.5,
                    size=22, bold=True, color=WHITE)

        # Key stats row
        stat_items = []
        if facts.get("companies"):
            stat_items.append((str(len(facts["companies"])), "Companies Analyzed"))
        if facts.get("how_it_works"):
            stat_items.append((str(len(facts["how_it_works"])), "How-It-Works Steps"))
        if facts.get("key_components"):
            stat_items.append((str(len(facts["key_components"])), "Key Components"))
        if facts.get("market_facts"):
            stat_items.append((str(len(facts["market_facts"])), "Market Facts"))
        if facts.get("key_findings"):
            stat_items.append((str(len(facts["key_findings"])), "Key Findings"))
        if facts.get("use_cases"):
            stat_items.append((str(len(facts["use_cases"])), "Use Cases"))

        for i, (num, label) in enumerate(stat_items[:4]):
            x = 0.6 + i * 3.1
            add_rect(s, x, 2.3, 2.8, 1.6, SURFACE)
            add_rect(s, x, 2.3, 2.8, 0.05, badge_color)
            add_textbox(s, num,   x+0.15, 2.45, 2.5, 0.7, size=32, bold=True, color=badge_color)
            add_textbox(s, label, x+0.15, 3.15, 2.5, 0.5, size=9,  color=MUTED)

        # Powered by line
        add_textbox(s, "Researched by AI Agent  ·  Web + ArXiv + OpenAI",
                    0.6, 4.3, 12.1, 0.4, size=10, color=MUTED)

        # Bottom branding
        add_textbox(s, "AI Research Agent",
                    0.6, 5.2, 6.0, 0.6, size=20, bold=True, color=badge_color)
        add_textbox(s, f"Generated on {date.today()}  ·  Powered by OpenAI + Tavily",
                    0.6, 5.9, 10.0, 0.4, size=9, color=MUTED)
        add_textbox(s, "github.com/Sajidhussain19/research-agent",
                    0.6, 6.4, 10.0, 0.4, size=9, color=MUTED)

    # ── Build deck based on type ──────────────────────────────────────────────
    if query_type == "concept":
        concept_slides()
    elif query_type == "market":
        market_slides()
    elif query_type == "research":
        research_slides()
    elif query_type == "mixed":
        concept_slides()
        market_slides()
    else:
        market_slides()

    conclusion_slide()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()